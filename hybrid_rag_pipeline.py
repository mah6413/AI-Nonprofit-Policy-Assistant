"""Grounded RAG pipeline for nonprofit document analysis and generation."""

import csv
import html
import math
import os
import re
from pathlib import Path
from typing import Any, Dict, List, Optional, Sequence
from urllib.parse import quote

import openpyxl
from docx import Document
from openai import OpenAI
from pptx import Presentation
from pypdf import PdfReader


try:
    from dotenv import load_dotenv
except ImportError:
    load_dotenv = None

if load_dotenv is not None:
    load_dotenv()

DEFAULT_EMBEDDING_MODEL = "text-embedding-3-small"
DEFAULT_CHAT_MODEL = "gpt-4.1"
PIPELINE_VERSION = "updated-hybrid-rag-2026-07-21"
HYBRID_SEMANTIC_WEIGHT = 0.68
HYBRID_KEYWORD_WEIGHT = 0.32
KEYWORD_STOPWORDS = {
    "a", "about", "above", "after", "again", "against", "all", "am", "an", "and",
    "any", "are", "as", "at", "be", "because", "been", "before", "being", "below",
    "between", "both", "but", "by", "can", "could", "did", "do", "does", "doing",
    "down", "during", "each", "few", "for", "from", "further", "had", "has", "have",
    "having", "he", "her", "here", "hers", "him", "his", "how", "i", "if", "in",
    "into", "is", "it", "its", "itself", "just", "me", "more", "most", "my", "no",
    "nor", "not", "now", "of", "off", "on", "once", "only", "or", "other", "our",
    "ours", "out", "over", "own", "same", "she", "should", "so", "some", "such",
    "than", "that", "the", "their", "theirs", "them", "then", "there", "these",
    "they", "this", "those", "through", "to", "too", "under", "until", "up", "very",
    "was", "we", "were", "what", "when", "where", "which", "while", "who", "whom",
    "why", "will", "with", "you", "your", "yours",
}

BASE_SYSTEM_PROMPT = """
You are a Nonprofit Fundraising, Strategy, Policy, Communications, Programs,
Compliance, and Community Engagement Assistant.

Your job is to answer questions and draft professional nonprofit documents using
the uploaded source materials. Preserve the organization's voice, audience, and
goals when they are present in the sources.

Rules:
- Ground factual claims, figures, dates, quotations, outcomes, and legal or policy
  statements in the retrieved context.
- Never invent donors, grant requirements, financial figures, program outcomes,
  legal obligations, quotations, citations, partnerships, or approvals.
- Clearly label placeholders such as [DONOR NAME] or [AMOUNT NEEDED].
- State when the evidence is incomplete, inconsistent, or unavailable.
- Protect confidential client, employee, donor, and case-management information.
- Do not claim to provide legal, tax, accounting, or investment advice.
- Use clear, inclusive, audience-appropriate nonprofit language.
""".strip()


# Each document type maps to its category and a practical default outline. Callers
# can override the outline by including explicit requirements in user_request.
DOCUMENT_CATALOG: Dict[str, Dict[str, Any]] = {
    # Fundraising & Development Documents
    "grant_proposal": {"category": "fundraising_development", "title": "Grant Proposal", "sections": ["Cover Information", "Executive Summary", "Statement of Need", "Program Description", "Goals and Measurable Objectives", "Activities and Timeline", "Organizational Capacity", "Evaluation", "Sustainability", "Budget Narrative", "Conclusion"]},
    "grant_budget": {"category": "fundraising_development", "title": "Grant Budget", "sections": ["Budget Assumptions", "Revenue", "Personnel", "Program Expenses", "Administrative Expenses", "Indirect Costs", "Total Request", "Budget Narrative"]},
    "grant_report": {"category": "fundraising_development", "title": "Grant Report", "sections": ["Grant Overview", "Activities Completed", "Progress Toward Outcomes", "People Served", "Challenges and Adaptations", "Financial Summary", "Lessons Learned", "Next Steps"]},
    "letter_of_inquiry": {"category": "fundraising_development", "title": "Letter of Inquiry", "sections": ["Opening and Funding Request", "Organizational Overview", "Need", "Proposed Work", "Expected Results", "Budget and Request", "Closing"]},
    "foundation_prospect_research": {"category": "fundraising_development", "title": "Foundation Prospect Research", "sections": ["Prospect Snapshot", "Funding Priorities", "Geographic and Population Fit", "Typical Grants", "Eligibility and Restrictions", "Relevant Relationships", "Alignment Assessment", "Recommended Approach", "Open Questions"]},
    "donor_stewardship_report": {"category": "fundraising_development", "title": "Donor Stewardship Report", "sections": ["Message of Appreciation", "Gift Overview", "Impact Highlights", "Program Results", "Financial Stewardship", "Stories or Examples", "Looking Ahead"]},
    "sponsorship_package": {"category": "fundraising_development", "title": "Sponsorship Package", "sections": ["Opportunity Overview", "Organization and Audience", "Event or Campaign Details", "Sponsorship Levels", "Benefits and Recognition", "Impact of Support", "Commitment Instructions"]},
    "case_statement": {"category": "fundraising_development", "title": "Fundraising Case Statement", "sections": ["Vision", "Urgent Need", "Our Response", "Evidence of Impact", "Campaign Priorities", "Funding Goal and Use of Funds", "Why Now", "Call to Partnership"]},

    # Donor & Development Communications
    "donation_appeal_letter": {"category": "donor_communications", "title": "Donation Appeal Letter", "sections": ["Personalized Opening", "Need or Story", "Organizational Response", "Evidence of Impact", "Specific Ask", "Call to Action", "Gratitude and Closing"]},
    "annual_giving_email": {"category": "donor_communications", "title": "Annual Giving Campaign Email", "sections": ["Subject Line", "Preview Text", "Opening", "Impact Message", "Specific Ask", "Call-to-Action Button Text", "Closing", "P.S."]},
    "major_donor_briefing": {"category": "donor_communications", "title": "Major Donor Briefing", "sections": ["Donor Profile", "Relationship History", "Interests and Alignment", "Recent Engagement", "Giving Context", "Meeting Objective", "Key Messages", "Suggested Ask", "Sensitive Topics and Follow-Up"]},
    "development_plan": {"category": "donor_communications", "title": "Development Plan", "sections": ["Executive Summary", "Fundraising Goals", "Revenue Mix", "Donor Segments", "Strategies and Tactics", "Calendar and Milestones", "Roles and Responsibilities", "Budget and Resources", "KPIs", "Risks and Contingencies"]},
    "capital_campaign_material": {"category": "donor_communications", "title": "Capital Campaign Material", "sections": ["Campaign Vision", "Case for Support", "Campaign Goal", "Funding Priorities", "Community Impact", "Gift Opportunities", "Campaign Timeline", "Call to Action"]},
    "donor_impact_report": {"category": "donor_communications", "title": "Donor Impact Report", "sections": ["Thank You", "Impact at a Glance", "How Funds Were Used", "Results and Outcomes", "Participant or Community Perspective", "Financial Snapshot", "What Comes Next"]},

    # Strategic & Executive Documents
    "strategic_plan": {"category": "strategic_executive", "title": "Strategic Plan", "sections": ["Executive Summary", "Mission, Vision, and Values", "Environmental Context", "Strategic Priorities", "Goals and Objectives", "Initiatives", "Measures of Success", "Implementation Timeline", "Accountability and Review"]},
    "annual_operating_plan": {"category": "strategic_executive", "title": "Annual Operating Plan", "sections": ["Annual Priorities", "Objectives and Deliverables", "Activities", "Owners", "Timeline", "Resources and Budget", "KPIs", "Dependencies and Risks", "Review Cadence"]},
    "swot_analysis": {"category": "strategic_executive", "title": "SWOT Analysis", "sections": ["Purpose and Scope", "Strengths", "Weaknesses", "Opportunities", "Threats", "Strategic Implications", "Priority Actions"]},
    "theory_of_change": {"category": "strategic_executive", "title": "Theory of Change", "sections": ["Long-Term Impact", "Problem and Context", "Priority Populations", "Preconditions", "Activities", "Outputs", "Short- and Medium-Term Outcomes", "Assumptions", "External Factors", "Indicators"]},
    "logic_model": {"category": "strategic_executive", "title": "Logic Model", "sections": ["Situation", "Inputs", "Activities", "Outputs", "Short-Term Outcomes", "Intermediate Outcomes", "Long-Term Outcomes", "Assumptions and External Factors", "Indicators"]},
    "kpi_dashboard": {"category": "strategic_executive", "title": "KPI Dashboard", "sections": ["Reporting Period", "Strategic Objective", "KPI", "Definition", "Baseline", "Target", "Current Result", "Status", "Trend", "Owner", "Corrective Action"]},
    "program_evaluation_report": {"category": "strategic_executive", "title": "Program Evaluation Report", "sections": ["Executive Summary", "Program Background", "Evaluation Questions", "Methodology", "Findings", "Equity Considerations", "Limitations", "Conclusions", "Recommendations"]},
    "board_agenda": {"category": "strategic_executive", "title": "Board Meeting Agenda", "sections": ["Meeting Details", "Call to Order", "Consent Agenda", "Reports", "Discussion Items", "Decision Items", "Executive Session", "Action Review", "Adjournment"]},
    "board_minutes": {"category": "strategic_executive", "title": "Board Minutes", "sections": ["Meeting Details", "Attendance and Quorum", "Approval of Prior Minutes", "Reports Received", "Discussion Summary", "Motions and Votes", "Conflicts or Recusals", "Action Items", "Adjournment"]},
    "board_packet": {"category": "strategic_executive", "title": "Board Packet", "sections": ["Cover and Meeting Details", "Agenda", "Prior Minutes", "Executive Director Report", "Financial Reports", "Committee Reports", "Decision Memos", "Background Materials", "Action Register"]},
    "bylaws": {"category": "strategic_executive", "title": "Bylaws Draft", "sections": ["Name and Purpose", "Members", "Board Composition", "Officers", "Meetings and Voting", "Committees", "Conflicts of Interest", "Indemnification", "Fiscal Administration", "Amendments", "Dissolution"]},
    "governance_policy": {"category": "strategic_executive", "title": "Governance Policy", "sections": ["Purpose", "Scope", "Definitions", "Policy", "Roles and Responsibilities", "Procedures", "Documentation", "Review and Approval"]},
    "executive_director_report": {"category": "strategic_executive", "title": "Executive Director Report", "sections": ["Executive Summary", "Mission and Program Updates", "Financial and Development Update", "People and Operations", "External Relations", "Risks and Decisions Needed", "Priorities Before Next Meeting"]},
    "committee_report": {"category": "strategic_executive", "title": "Committee Report", "sections": ["Committee and Reporting Period", "Members", "Work Completed", "Key Findings", "Recommendations", "Decisions Requested", "Next Steps"]},

    # Policy, Advocacy & Research Documents
    "policy_brief": {"category": "policy_advocacy_research", "title": "Policy Brief", "sections": ["Executive Summary", "Issue Overview", "Evidence and Findings", "Policy Context", "Options", "Recommendations", "Conclusion"]},
    "legislative_memo": {"category": "policy_advocacy_research", "title": "Legislative Memo", "sections": ["To, From, Date, and Re", "Bottom Line", "Bill or Issue Summary", "Current Status", "Analysis", "Stakeholder Positions", "Implications", "Recommended Action"]},
    "advocacy_letter": {"category": "policy_advocacy_research", "title": "Advocacy Letter", "sections": ["Date and Recipient", "Position", "Community Impact", "Supporting Evidence", "Specific Policy Ask", "Closing"]},
    "public_comment": {"category": "policy_advocacy_research", "title": "Public Comment", "sections": ["Agency and Docket", "Commenter Introduction", "Summary of Position", "Legal or Policy Context", "Evidence and Analysis", "Detailed Recommendations", "Conclusion"]},
    "white_paper": {"category": "policy_advocacy_research", "title": "White Paper", "sections": ["Executive Summary", "Problem Definition", "Background", "Methodology", "Analysis", "Findings", "Options", "Recommendations", "Conclusion"]},
    "regulatory_analysis": {"category": "policy_advocacy_research", "title": "Regulatory Analysis", "sections": ["Rule Overview", "Authority and Context", "Key Provisions", "Affected Stakeholders", "Operational and Equity Impacts", "Compliance Considerations", "Ambiguities and Risks", "Recommended Response"]},
    "hearing_testimony": {"category": "policy_advocacy_research", "title": "Hearing Testimony", "sections": ["Chair and Committee Salutation", "Witness Introduction", "Position", "Evidence and Community Impact", "Recommendations", "Closing and Offer to Answer Questions"]},
    "issue_brief": {"category": "policy_advocacy_research", "title": "Issue Brief", "sections": ["Issue at a Glance", "Background", "Who Is Affected", "Key Evidence", "Current Landscape", "Considerations", "Next Steps"]},
    "research_report": {"category": "policy_advocacy_research", "title": "Research Report", "sections": ["Abstract", "Introduction", "Research Questions", "Methods", "Findings", "Discussion", "Limitations", "Recommendations", "References"]},
    "stakeholder_analysis": {"category": "policy_advocacy_research", "title": "Stakeholder Analysis", "sections": ["Issue and Objective", "Stakeholder", "Interests", "Influence", "Position", "Relationships", "Risks", "Engagement Strategy"]},
    "campaign_strategy_memo": {"category": "policy_advocacy_research", "title": "Campaign Strategy Memo", "sections": ["Objective", "Landscape", "Targets", "Audiences", "Theory of Influence", "Messages", "Tactics and Timeline", "Partners", "Metrics", "Risks"]},
    "grassroots_mobilization_plan": {"category": "policy_advocacy_research", "title": "Grassroots Mobilization Plan", "sections": ["Goal", "Target Audiences", "Recruitment", "Leadership Development", "Mobilization Tactics", "Messaging", "Calendar", "Data and Tools", "Metrics", "Risk and Safety"]},
    "legislative_tracker": {"category": "policy_advocacy_research", "title": "Legislative Tracking Sheet", "sections": ["Bill Number", "Title", "Sponsor", "Summary", "Status", "Last Action", "Next Deadline", "Organizational Position", "Priority", "Owner", "Notes"]},
    "talking_points": {"category": "policy_advocacy_research", "title": "Talking Points", "sections": ["Core Message", "Three Key Points", "Supporting Evidence", "Audience Relevance", "Anticipated Questions", "Response Guidance", "Call to Action"]},
    "fact_sheet": {"category": "policy_advocacy_research", "title": "Fact Sheet", "sections": ["Headline", "Issue in Brief", "Key Facts", "Who Is Affected", "What Works", "Organizational Response", "Call to Action", "Sources"]},

    # Communications & Marketing Documents
    "press_release": {"category": "communications_marketing", "title": "Press Release", "sections": ["Release Status and Date", "Headline", "Subheadline", "Dateline and Lead", "Key Details", "Attributed Quote Placeholders", "Organization Boilerplate", "Media Contact"]},
    "newsletter": {"category": "communications_marketing", "title": "Newsletter", "sections": ["Subject Line", "Opening Note", "Lead Story", "Program Updates", "Impact Highlight", "Upcoming Dates", "Calls to Action", "Closing"]},
    "social_media_calendar": {"category": "communications_marketing", "title": "Social Media Calendar", "sections": ["Date", "Platform", "Campaign", "Audience", "Post Copy", "Creative Direction", "Link", "Call to Action", "Owner", "Status"]},
    "media_kit": {"category": "communications_marketing", "title": "Media Kit", "sections": ["Organization Overview", "Key Facts", "Leadership Bios", "Programs and Impact", "Approved Messages", "Brand Assets", "Recent Coverage", "Media Contact"]},
    "website_content": {"category": "communications_marketing", "title": "Website Content Draft", "sections": ["Page Goal", "Audience", "Page Title", "Hero Message", "Body Copy", "Proof Points", "Calls to Action", "SEO Title and Description"]},
    "email_campaign": {"category": "communications_marketing", "title": "Email Campaign", "sections": ["Campaign Goal", "Audience Segment", "Email Sequence", "Subject Lines", "Preview Text", "Message Copy", "Calls to Action", "Send Timing", "Success Metrics"]},
    "annual_report": {"category": "communications_marketing", "title": "Annual Report", "sections": ["Leadership Message", "Year at a Glance", "Mission and Strategy", "Program Highlights", "Outcomes and Impact", "Community Voices", "Financial Summary", "Supporters", "Looking Ahead"]},
    "messaging_guide": {"category": "communications_marketing", "title": "Messaging Guide", "sections": ["Purpose and Audiences", "Core Narrative", "Message Pillars", "Proof Points", "Audience-Specific Messages", "Preferred Language", "Language to Avoid", "Elevator Pitch", "FAQs"]},
    "communication_strategy": {"category": "communications_marketing", "title": "Communication Strategy", "sections": ["Situation Analysis", "Objectives", "Audiences", "Messages", "Channels and Tactics", "Content Calendar", "Roles and Workflow", "Budget", "Measurement", "Risks"]},
    "brand_guidelines": {"category": "communications_marketing", "title": "Brand Guidelines", "sections": ["Brand Foundation", "Voice and Tone", "Logo Use", "Color Palette", "Typography", "Imagery", "Accessibility", "Templates", "Governance"]},
    "speech_draft": {"category": "communications_marketing", "title": "Speech Draft", "sections": ["Opening", "Audience Connection", "Core Message", "Supporting Story or Evidence", "Key Points", "Call to Action", "Closing"]},
    "event_script": {"category": "communications_marketing", "title": "Event Script", "sections": ["Event Details", "Opening and Welcome", "Transitions", "Speaker Introductions", "Program Segments", "Recognition", "Call to Action", "Closing", "Production Cues"]},

    # Program & Service Delivery Documents
    "program_manual": {"category": "program_service_delivery", "title": "Program Manual", "sections": ["Purpose and Scope", "Program Model", "Eligibility", "Roles", "Service Workflow", "Documentation", "Quality Standards", "Safety and Escalation", "Data and Confidentiality", "Continuous Improvement"]},
    "client_intake_form": {"category": "program_service_delivery", "title": "Client Intake Form", "sections": ["Client Information", "Contact Preferences", "Demographics", "Needs and Goals", "Eligibility", "Accessibility and Language", "Emergency or Safety Information", "Consent and Privacy", "Staff Use"]},
    "case_management_note": {"category": "program_service_delivery", "title": "Case Management Note", "sections": ["Client Identifier", "Date and Contact Type", "Purpose", "Objective Observations", "Services and Referrals", "Client Response", "Risk or Safety", "Plan and Follow-Up", "Staff Name"]},
    "training_material": {"category": "program_service_delivery", "title": "Training Material", "sections": ["Learning Objectives", "Audience and Prerequisites", "Agenda", "Core Content", "Practice Activity", "Scenario or Example", "Knowledge Check", "Resources", "Evaluation"]},
    "curriculum": {"category": "program_service_delivery", "title": "Curriculum", "sections": ["Course Overview", "Learner Profile", "Learning Outcomes", "Scope and Sequence", "Lesson Plans", "Activities", "Materials", "Assessment", "Accessibility and Adaptation"]},
    "standard_operating_procedure": {"category": "program_service_delivery", "title": "Standard Operating Procedure", "sections": ["Purpose", "Scope", "Definitions", "Roles", "Required Materials or Systems", "Procedure", "Exceptions and Escalation", "Records", "Quality Control", "Revision History"]},
    "volunteer_handbook": {"category": "program_service_delivery", "title": "Volunteer Handbook", "sections": ["Welcome and Mission", "Volunteer Roles", "Onboarding", "Code of Conduct", "Safeguarding and Confidentiality", "Attendance and Communication", "Safety", "Problem Resolution", "Recognition", "Acknowledgment"]},
    "survey_results": {"category": "program_service_delivery", "title": "Survey Results Report", "sections": ["Purpose", "Respondents", "Method", "Key Results", "Subgroup Findings", "Open-Ended Themes", "Limitations", "Implications", "Recommended Actions"]},
    "outcome_measurement_report": {"category": "program_service_delivery", "title": "Outcome Measurement Report", "sections": ["Outcome Framework", "Indicators", "Data Sources", "Baseline and Targets", "Results", "Interpretation", "Equity Analysis", "Limitations", "Improvement Actions"]},
    "impact_assessment": {"category": "program_service_delivery", "title": "Impact Assessment", "sections": ["Executive Summary", "Intervention and Context", "Assessment Design", "Evidence", "Outcomes and Impact", "Unintended Effects", "Equity Considerations", "Limitations", "Recommendations"]},
    "program_metrics": {"category": "program_service_delivery", "title": "Program Performance Metrics", "sections": ["Reporting Period", "Metric", "Definition", "Data Source", "Baseline", "Target", "Actual", "Variance", "Interpretation", "Action"]},

    # Financial & Compliance Documents
    "budget": {"category": "financial_compliance", "title": "Budget", "sections": ["Assumptions", "Revenue", "Personnel", "Program Expenses", "Administrative Expenses", "Fundraising Expenses", "Capital Expenses", "Total and Variance Notes"]},
    "financial_statement": {"category": "financial_compliance", "title": "Financial Statement Narrative", "sections": ["Reporting Period", "Statement of Financial Position", "Statement of Activities", "Cash Flow", "Budget-to-Actual", "Restrictions and Liquidity", "Material Variances", "Management Notes"]},
    "audit_report_summary": {"category": "financial_compliance", "title": "Audit Report Summary", "sections": ["Audit Scope and Period", "Auditor Opinion", "Key Financial Findings", "Internal Control Findings", "Compliance Findings", "Management Responses", "Required Actions"]},
    "expense_tracker": {"category": "financial_compliance", "title": "Expense Tracking Sheet", "sections": ["Date", "Vendor", "Description", "Program or Department", "Account", "Funding Source", "Amount", "Approval", "Payment Status", "Receipt"]},
    "form_990_support": {"category": "financial_compliance", "title": "IRS Form 990 Documentation Checklist", "sections": ["Organization Information", "Governance", "Programs and Accomplishments", "Revenue", "Expenses", "Compensation", "Related Organizations and Transactions", "Policies", "Schedules and Supporting Records", "Review and Approval"]},
    "forecasting_model": {"category": "financial_compliance", "title": "Forecasting Model Narrative", "sections": ["Forecast Period", "Assumptions", "Revenue Scenarios", "Expense Drivers", "Cash Projection", "Restricted Funds", "Scenario Analysis", "Risks", "Management Actions"]},
    "employee_handbook": {"category": "financial_compliance", "title": "Employee Handbook", "sections": ["Welcome and Mission", "Employment Practices", "Compensation and Timekeeping", "Benefits and Leave", "Workplace Conduct", "Performance", "Safety and Security", "Technology and Data", "Separation", "Acknowledgment"]},
    "hr_policy": {"category": "financial_compliance", "title": "HR Policy", "sections": ["Purpose", "Scope", "Definitions", "Policy", "Responsibilities", "Procedure", "Reporting and Non-Retaliation", "Records", "Review"]},
    "dei_policy": {"category": "financial_compliance", "title": "DEI Policy", "sections": ["Commitment", "Purpose and Scope", "Definitions", "Standards", "Accessibility and Inclusion", "Accountability", "Reporting Concerns", "Measurement and Review"]},
    "cybersecurity_policy": {"category": "financial_compliance", "title": "Cybersecurity Policy", "sections": ["Purpose and Scope", "Roles", "Access Control", "Devices and Networks", "Data Classification", "Backups", "Vendor Security", "Incident Response", "Training", "Review"]},
    "risk_management_plan": {"category": "financial_compliance", "title": "Risk Management Plan", "sections": ["Context and Risk Appetite", "Risk Register", "Likelihood and Impact", "Controls", "Owners", "Response Plans", "Monitoring", "Incident Escalation", "Board Reporting"]},
    "vendor_contract": {"category": "financial_compliance", "title": "Vendor Contract Draft", "sections": ["Parties", "Scope and Deliverables", "Term", "Fees and Payment", "Responsibilities", "Confidentiality and Data", "Insurance and Compliance", "Termination", "Dispute Terms", "Signatures"]},
    "memorandum_of_understanding": {"category": "financial_compliance", "title": "Memorandum of Understanding", "sections": ["Parties and Purpose", "Shared Goals", "Roles and Contributions", "Governance and Communication", "Data and Confidentiality", "Funding and Resources", "Term and Modification", "Dispute Resolution", "Signatures"]},

    # Event & Community Engagement Documents
    "event_run_of_show": {"category": "event_community_engagement", "title": "Event Run of Show", "sections": ["Event Details", "Time", "Duration", "Segment", "Speaker or Owner", "Script or Action", "Production Cue", "Materials", "Contingency"]},
    "speaker_bio": {"category": "event_community_engagement", "title": "Speaker Bio", "sections": ["Name and Role", "Current Work", "Relevant Expertise", "Selected Accomplishments", "Connection to Event Topic", "Pronunciation and Pronouns"]},
    "conference_agenda": {"category": "event_community_engagement", "title": "Conference Agenda", "sections": ["Conference Details", "Registration", "Opening", "Sessions", "Breaks", "Networking", "Closing", "Accessibility and Logistics"]},
    "registration_tracker": {"category": "event_community_engagement", "title": "Registration Tracking Sheet", "sections": ["Registrant", "Organization", "Contact", "Ticket Type", "Registration Date", "Payment", "Accessibility", "Dietary Needs", "Consent", "Check-In"]},
    "sponsorship_agreement": {"category": "event_community_engagement", "title": "Sponsorship Agreement", "sections": ["Parties and Event", "Sponsorship Level", "Sponsor Contribution", "Organizer Benefits", "Brand and Recognition Terms", "Deadlines", "Cancellation", "Compliance", "Signatures"]},
    "community_needs_assessment": {"category": "event_community_engagement", "title": "Community Needs Assessment", "sections": ["Executive Summary", "Community Profile", "Purpose and Questions", "Methods", "Assets and Strengths", "Needs and Gaps", "Priority Populations", "Findings", "Recommendations", "Limitations"]},
    "stakeholder_engagement_report": {"category": "event_community_engagement", "title": "Stakeholder Engagement Report", "sections": ["Purpose", "Stakeholders Engaged", "Methods", "Participation", "Themes", "Areas of Agreement and Tension", "How Input Will Be Used", "Follow-Up Commitments"]},
    "outreach_script": {"category": "event_community_engagement", "title": "Outreach Script", "sections": ["Introduction", "Reason for Contact", "Audience-Relevant Message", "Key Information", "Conversation Questions", "Call to Action", "Objection Responses", "Closing and Follow-Up"]},
    "volunteer_coordination_plan": {"category": "event_community_engagement", "title": "Volunteer Coordination Plan", "sections": ["Objective", "Roles and Staffing", "Recruitment", "Screening and Onboarding", "Training", "Scheduling", "Communication", "Safety and Escalation", "Recognition", "Evaluation"]},
}

CATEGORY_TITLES = {
    "fundraising_development": "Fundraising & Development Documents",
    "donor_communications": "Donor & Development Communications",
    "strategic_executive": "Strategic & Executive Documents",
    "policy_advocacy_research": "Policy, Advocacy & Research Documents",
    "communications_marketing": "Communications & Marketing Documents",
    "program_service_delivery": "Program & Service Delivery Documents",
    "financial_compliance": "Financial & Compliance Documents",
    "event_community_engagement": "Event & Community Engagement Documents",
}


def get_openai_client(api_key: Optional[str] = None) -> OpenAI:
    """Create a client lazily so importing this module does not require a key."""
    key = api_key or os.getenv("OPENAI_API_KEY")
    if not key:
        raise ValueError("OPENAI_API_KEY is not set in the environment or .env file.")
    return OpenAI(api_key=key)


def clean_text(text: str) -> str:
    return re.sub(r"\s+", " ", text.replace("\x00", " ")).strip()


def chunk_text(text: str, chunk_size: int = 1200, overlap: int = 200) -> List[str]:
    if chunk_size <= 0 or overlap < 0 or overlap >= chunk_size:
        raise ValueError("chunk_size must be positive and overlap must be smaller.")
    text = clean_text(text)
    step = chunk_size - overlap
    return [text[start:start + chunk_size] for start in range(0, len(text), step) if text[start:start + chunk_size].strip()]


def cosine_similarity(a: Sequence[float], b: Sequence[float]) -> float:
    if len(a) != len(b):
        raise ValueError("Embedding vectors must have the same length.")
    norm_a = math.sqrt(sum(x * x for x in a))
    norm_b = math.sqrt(sum(x * x for x in b))
    return 0.0 if not norm_a or not norm_b else sum(x * y for x, y in zip(a, b)) / (norm_a * norm_b)


def tokenize_for_search(text: str) -> List[str]:
    """Normalize text into lightweight keyword tokens for hybrid retrieval."""
    tokens = re.findall(r"[a-z0-9][a-z0-9'-]*", text.lower())
    return [token for token in tokens if len(token) > 2 and token not in KEYWORD_STOPWORDS]


def keyword_overlap_score(query: str, text: str) -> float:
    """Return a bounded keyword relevance score using term overlap and coverage."""
    query_terms = tokenize_for_search(query)
    if not query_terms:
        return 0.0
    text_terms = tokenize_for_search(text)
    if not text_terms:
        return 0.0

    query_counts: Dict[str, int] = {}
    for term in query_terms:
        query_counts[term] = query_counts.get(term, 0) + 1

    text_counts: Dict[str, int] = {}
    for term in text_terms:
        text_counts[term] = text_counts.get(term, 0) + 1

    overlap = 0.0
    matched_terms = 0
    for term, query_count in query_counts.items():
        text_count = text_counts.get(term, 0)
        if text_count:
            matched_terms += 1
            overlap += min(query_count, text_count) * (1.0 + math.log1p(text_count))

    coverage = matched_terms / max(1, len(query_counts))
    density = overlap / max(1.0, len(text_terms) ** 0.5)
    return min(1.0, (0.65 * coverage) + (0.35 * min(1.0, density)))


def clone_chunk_for_result(chunk: Dict[str, Any]) -> Dict[str, Any]:
    return {
        key: chunk[key]
        for key in (
            "chunk_id",
            "document_id",
            "document_title",
            "document_type",
            "source_path",
            "location",
            "locator",
            "char_start",
            "char_end",
            "text",
        )
    }


def parse_pdf(file_path: str) -> str:
    reader = PdfReader(file_path)
    return "\n".join(f"[Page {index}]\n{page.extract_text() or ''}" for index, page in enumerate(reader.pages, 1))


def parse_docx(file_path: str) -> str:
    doc = Document(file_path)
    blocks = [p.text for p in doc.paragraphs if p.text.strip()]
    for table_index, table in enumerate(doc.tables, 1):
        blocks.append(f"[Table {table_index}]")
        blocks.extend(" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows)
    return "\n".join(blocks)


def parse_pptx(file_path: str) -> str:
    presentation = Presentation(file_path)
    slides = []
    for index, slide in enumerate(presentation.slides, 1):
        text = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        slides.append(f"[Slide {index}]\n" + "\n".join(text))
    return "\n".join(slides)


def parse_xlsx(file_path: str) -> str:
    workbook = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
    blocks: List[str] = []
    for sheet in workbook.worksheets:
        blocks.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            values = [str(value) for value in row if value is not None]
            if values:
                blocks.append(" | ".join(values))
    return "\n".join(blocks)


def parse_delimited(file_path: str) -> str:
    delimiter = "\t" if file_path.lower().endswith(".tsv") else ","
    with open(file_path, "r", encoding="utf-8-sig", newline="") as source:
        return "\n".join(" | ".join(row) for row in csv.reader(source, delimiter=delimiter))


def parse_file_segments(file_path: str) -> List[Dict[str, Any]]:
    """Extract independently citable units with document-native locations."""
    extension = Path(file_path).suffix.lower()
    segments: List[Dict[str, Any]] = []

    if extension == ".pdf":
        for page_number, page in enumerate(PdfReader(file_path).pages, 1):
            segments.append({"text": page.extract_text() or "", "location": f"Page {page_number}", "locator": {"page": page_number}})
    elif extension == ".docx":
        doc = Document(file_path)
        for paragraph_number, paragraph in enumerate(doc.paragraphs, 1):
            if paragraph.text.strip():
                segments.append({"text": paragraph.text, "location": f"Paragraph {paragraph_number}", "locator": {"paragraph": paragraph_number}})
        for table_number, table in enumerate(doc.tables, 1):
            for row_number, row in enumerate(table.rows, 1):
                text = " | ".join(cell.text.strip() for cell in row.cells)
                if text.strip(" |"):
                    segments.append({"text": text, "location": f"Table {table_number}, row {row_number}", "locator": {"table": table_number, "row": row_number}})
    elif extension == ".pptx":
        presentation = Presentation(file_path)
        for slide_number, slide in enumerate(presentation.slides, 1):
            text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip())
            if text.strip():
                segments.append({"text": text, "location": f"Slide {slide_number}", "locator": {"slide": slide_number}})
    elif extension == ".xlsx":
        workbook = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
        for sheet in workbook.worksheets:
            for row_number, row in enumerate(sheet.iter_rows(values_only=True), 1):
                values = [str(value) for value in row if value is not None]
                if values:
                    end_column = openpyxl.utils.get_column_letter(max(1, len(row)))
                    cell_range = f"A{row_number}:{end_column}{row_number}"
                    segments.append({"text": " | ".join(values), "location": f"Sheet '{sheet.title}', cells {cell_range}", "locator": {"sheet": sheet.title, "cell_range": cell_range}})
    elif extension in {".csv", ".tsv"}:
        delimiter = "\t" if extension == ".tsv" else ","
        with open(file_path, "r", encoding="utf-8-sig", newline="") as source:
            for row_number, row in enumerate(csv.reader(source, delimiter=delimiter), 1):
                if any(cell.strip() for cell in row):
                    segments.append({"text": " | ".join(row), "location": f"Row {row_number}", "locator": {"row": row_number}})
    elif extension in {".txt", ".md"}:
        for line_number, line in enumerate(Path(file_path).read_text(encoding="utf-8-sig").splitlines(), 1):
            if line.strip():
                segments.append({"text": line, "location": f"Line {line_number}", "locator": {"line": line_number}})
    else:
        raise ValueError(f"Unsupported file format: {extension or '[none]'}")

    return segments


def parse_uploaded_file(file_path: str) -> Dict[str, str]:
    extension = Path(file_path).suffix.lower()
    parsers = {".pdf": parse_pdf, ".docx": parse_docx, ".pptx": parse_pptx, ".xlsx": parse_xlsx, ".csv": parse_delimited, ".tsv": parse_delimited}
    if extension in {".txt", ".md"}:
        return {"document_type": extension[1:], "text": Path(file_path).read_text(encoding="utf-8-sig")}
    if extension not in parsers:
        raise ValueError(f"Unsupported file format: {extension or '[none]'}")
    return {"document_type": extension[1:], "text": parsers[extension](file_path)}


def build_document_library(file_paths: List[str]) -> List[Dict[str, Any]]:
    library = []
    for index, file_path in enumerate(file_paths, 1):
        parsed = parse_uploaded_file(file_path)
        library.append({"document_id": f"doc_{index}", "document_title": Path(file_path).name, "document_type": parsed["document_type"], "source_path": str(Path(file_path).resolve()), "raw_text": parsed["text"], "segments": parse_file_segments(file_path)})
    return library


def build_chunk_store(document_library: List[Dict[str, Any]], chunk_size: int = 1200, overlap: int = 200) -> List[Dict[str, Any]]:
    store = []
    for document in document_library:
        chunk_index = 0
        for segment in document.get("segments") or [{"text": document["raw_text"], "location": "Document", "locator": {}}]:
            segment_text = clean_text(segment["text"])
            step = chunk_size - overlap
            for start in range(0, len(segment_text), step):
                text = segment_text[start:start + chunk_size]
                if not text.strip():
                    continue
                chunk_index += 1
                store.append({"chunk_id": f"{document['document_id']}_chunk_{chunk_index}", "document_id": document["document_id"], "document_title": document["document_title"], "document_type": document["document_type"], "source_path": document["source_path"], "location": segment["location"], "locator": segment["locator"], "char_start": start, "char_end": start + len(text), "text": text, "embedding": None})
    return store


def get_embedding(text: str, client: OpenAI, model: str = DEFAULT_EMBEDDING_MODEL) -> List[float]:
    return client.embeddings.create(model=model, input=text).data[0].embedding


def embed_chunk_store(chunk_store: List[Dict[str, Any]], client: OpenAI, embedding_model: str = DEFAULT_EMBEDDING_MODEL) -> List[Dict[str, Any]]:
    for chunk in chunk_store:
        if chunk["embedding"] is None:
            chunk["embedding"] = get_embedding(chunk["text"], client, embedding_model)
    return chunk_store


def retrieve_relevant_chunks(
    question: str,
    chunk_store: List[Dict[str, Any]],
    client: OpenAI,
    top_k: int = 6,
    embedding_model: str = DEFAULT_EMBEDDING_MODEL,
    semantic_weight: float = HYBRID_SEMANTIC_WEIGHT,
    keyword_weight: float = HYBRID_KEYWORD_WEIGHT,
) -> List[Dict[str, Any]]:
    """Retrieve with hybrid RAG: semantic embeddings plus keyword matching.

    The semantic side catches meaning even when the wording differs. The keyword
    side keeps exact names, grant terms, bill numbers, dates, and budget labels
    from getting washed out by embedding similarity alone.
    """
    if top_k <= 0:
        return []

    question_embedding = get_embedding(question, client, embedding_model)
    total_weight = semantic_weight + keyword_weight
    semantic_ratio = semantic_weight / total_weight if total_weight else 0.5
    keyword_ratio = keyword_weight / total_weight if total_weight else 0.5
    scored: List[Dict[str, Any]] = []

    for chunk in chunk_store:
        semantic_score = 0.0
        if chunk.get("embedding") is not None:
            semantic_score = cosine_similarity(question_embedding, chunk["embedding"])
            semantic_score = max(0.0, min(1.0, semantic_score))
        keyword_score = keyword_overlap_score(question, chunk["text"])
        hybrid_score = (semantic_ratio * semantic_score) + (keyword_ratio * keyword_score)

        if hybrid_score <= 0:
            continue
        item = clone_chunk_for_result(chunk)
        item["score"] = hybrid_score
        item["semantic_score"] = semantic_score
        item["keyword_score"] = keyword_score
        item["retrieval_method"] = "hybrid"
        scored.append(item)

    scored.sort(
        key=lambda item: (
            item["score"],
            item["keyword_score"],
            item["semantic_score"],
        ),
        reverse=True,
    )
    return scored[:top_k]


def evaluate_retrieval_quality(retrieved_chunks: List[Dict[str, Any]]) -> str:
    if not retrieved_chunks or retrieved_chunks[0]["score"] < 0.24:
        return "no_relevant_context"
    return "partial" if retrieved_chunks[0]["score"] < 0.50 else "strong"


def build_citations(chunks: List[Dict[str, Any]], citation_base_url: Optional[str] = None) -> List[Dict[str, Any]]:
    """Create stable citation records suitable for an in-app source viewer."""
    citations = []
    for index, item in enumerate(chunks, 1):
        citation_id = f"C{index}"
        source_url = Path(item["source_path"]).as_uri()
        if item["document_type"] == "pdf" and item["locator"].get("page"):
            source_url += f"#page={item['locator']['page']}"
        if citation_base_url:
            encoded_chunk_id = quote(item["chunk_id"])
            url = citation_base_url.format(chunk_id=encoded_chunk_id) if "{chunk_id}" in citation_base_url else f"{citation_base_url.rstrip('/')}/{encoded_chunk_id}"
        else:
            url = source_url
        citations.append({
            "id": citation_id,
            "chunk_id": item["chunk_id"],
            "url": url,
            "source_url": source_url,
            "document_id": item["document_id"],
            "document_title": item["document_title"],
            "document_type": item["document_type"],
            "location": item["location"],
            "locator": item["locator"],
            "char_start": item["char_start"],
            "char_end": item["char_end"],
            "excerpt": item["text"],
            "relevance": round(item["score"], 4),
            "semantic_relevance": round(item.get("semantic_score", 0.0), 4),
            "keyword_relevance": round(item.get("keyword_score", 0.0), 4),
            "retrieval_method": item.get("retrieval_method", "hybrid"),
        })
    return citations


def _context_from_chunks(chunks: List[Dict[str, Any]]) -> str:
    return "\n\n".join(
        (
            f"[C{index} | {item['document_title']} | {item['location']} | "
            f"hybrid={item['score']:.3f} | semantic={item.get('semantic_score', 0.0):.3f} | "
            f"keyword={item.get('keyword_score', 0.0):.3f}]\n{item['text']}"
        )
        for index, item in enumerate(chunks, 1)
    )


def link_citation_tokens(text: str, citations: List[Dict[str, Any]]) -> str:
    """Convert model tokens such as [C1] into clickable Markdown citations."""
    urls = {citation["id"]: citation["url"] for citation in citations}
    return re.sub(r"\[(C\d+)\]", lambda match: f"[{match.group(1)}]({urls[match.group(1)]})" if match.group(1) in urls else match.group(0), text)


def answer_question(question: str, chunk_store: List[Dict[str, Any]], client: OpenAI, top_k: int = 6, chat_model: str = DEFAULT_CHAT_MODEL, citation_base_url: Optional[str] = None) -> Dict[str, Any]:
    retrieved = retrieve_relevant_chunks(question, chunk_store, client, top_k)
    quality = evaluate_retrieval_quality(retrieved)
    if quality == "no_relevant_context":
        return {"answer": "I could not find enough support for that answer in the uploaded document repository.", "citations": [], "confidence": "low", "fallback_used": True, "retrieval_method": "hybrid", "pipeline_version": PIPELINE_VERSION}
    citations = build_citations(retrieved, citation_base_url)
    prompt = f"Question:\n{question}\n\nRetrieved context:\n{_context_from_chunks(retrieved)}\n\nAnswer from the context. Add the applicable citation token, such as [C1], immediately after every factual claim. Use only citation tokens present in the context. Distinguish sourced facts from gaps. Do not add a separate sources list."
    response = client.chat.completions.create(model=chat_model, messages=[{"role": "system", "content": BASE_SYSTEM_PROMPT}, {"role": "user", "content": prompt}])
    answer = link_citation_tokens(response.choices[0].message.content, citations)
    return {"answer": answer, "citations": citations, "confidence": "high" if quality == "strong" else "medium", "fallback_used": False, "retrieval_method": "hybrid", "pipeline_version": PIPELINE_VERSION}


def normalize_document_type(document_type: str) -> str:
    key = re.sub(r"[^a-z0-9]+", "_", document_type.lower()).strip("_")
    aliases = {
        "grant_proposals": "grant_proposal", "grant_application": "grant_proposal", "grant_applications": "grant_proposal",
        "loi": "letter_of_inquiry", "lois": "letter_of_inquiry", "grant_reporting_document": "grant_report",
        "sop": "standard_operating_procedure", "sops": "standard_operating_procedure", "mou": "memorandum_of_understanding", "mous": "memorandum_of_understanding",
        "board_meeting_agenda": "board_agenda", "board_meeting_minutes": "board_minutes", "kpi_dashboard": "kpi_dashboard",
        "annual_giving_campaign_email": "annual_giving_email", "capital_campaign_materials": "capital_campaign_material",
        "public_comments": "public_comment", "testimony_for_hearings": "hearing_testimony", "legislative_tracking_spreadsheet": "legislative_tracker",
        "case_statement_for_fundraising_campaign": "case_statement", "irs_form_990_documentation": "form_990_support",
        "mou_memorandum_of_understanding": "memorandum_of_understanding", "event_run_of_show_document": "event_run_of_show",
    }
    key = aliases.get(key, key)
    if key.endswith("s") and key[:-1] in DOCUMENT_CATALOG:
        key = key[:-1]
    if key not in DOCUMENT_CATALOG:
        raise ValueError(f"Unsupported document type: {document_type}. Use list_document_types() to see valid values.")
    return key


def list_document_types(category: Optional[str] = None) -> Dict[str, List[str]]:
    """Return supported document type keys grouped by category."""
    result: Dict[str, List[str]] = {}
    for key, spec in DOCUMENT_CATALOG.items():
        if category and spec["category"] != category:
            continue
        result.setdefault(spec["category"], []).append(key)
    return result


def generate_grounded_document(user_request: str, document_type: str, chunk_store: List[Dict[str, Any]], client: OpenAI, top_k: int = 8, chat_model: str = DEFAULT_CHAT_MODEL, citation_base_url: Optional[str] = None) -> Dict[str, Any]:
    key = normalize_document_type(document_type)
    spec = DOCUMENT_CATALOG[key]
    retrieval_query = f"{spec['title']}: {user_request}"
    retrieved = retrieve_relevant_chunks(retrieval_query, chunk_store, client, top_k)
    quality = evaluate_retrieval_quality(retrieved)
    if quality == "no_relevant_context":
        return {"answer": f"I could not find enough support in the uploaded repository to generate a reliable {spec['title']}.", "document_type": key, "category": spec["category"], "citations": [], "confidence": "low", "fallback_used": True, "retrieval_method": "hybrid", "pipeline_version": PIPELINE_VERSION}

    citations = build_citations(retrieved, citation_base_url)
    outline = "\n".join(f"- {section}" for section in spec["sections"])
    prompt = f"""
Create a {spec['title']} in the category {CATEGORY_TITLES[spec['category']]}.

User request:
{user_request}

Default structure:
{outline}

Retrieved source context:
{_context_from_chunks(retrieved)}

Requirements:
- Follow the requested format when it conflicts with the default structure.
- Use only supported facts and figures. Never turn assumptions into facts.
- Insert conspicuous [PLACEHOLDER] fields for necessary missing information.
- For table-oriented documents, produce a clean Markdown table.
- Flag substantive evidence gaps in a final "Information Needed" section.
- Place an applicable citation token such as [C1] immediately after each factual
  claim, figure, quotation, date, or sourced conclusion.
- Use only citation tokens supplied in the context and do not add a sources list.
- Draft legal, tax, HR, governance, contract, and compliance materials for qualified professional review.
""".strip()
    response = client.chat.completions.create(model=chat_model, messages=[{"role": "system", "content": BASE_SYSTEM_PROMPT}, {"role": "user", "content": prompt}])
    answer = link_citation_tokens(response.choices[0].message.content, citations)
    return {"answer": answer, "document_type": key, "category": spec["category"], "citations": citations, "confidence": "high" if quality == "strong" else "medium", "fallback_used": False, "retrieval_method": "hybrid", "pipeline_version": PIPELINE_VERSION}


def generate_fundraising_document(request: str, document_type: str, chunk_store: List[Dict[str, Any]], client: OpenAI, **kwargs: Any) -> Dict[str, Any]:
    return _generate_for_category("fundraising_development", request, document_type, chunk_store, client, **kwargs)


def generate_donor_communication(request: str, document_type: str, chunk_store: List[Dict[str, Any]], client: OpenAI, **kwargs: Any) -> Dict[str, Any]:
    return _generate_for_category("donor_communications", request, document_type, chunk_store, client, **kwargs)


def generate_strategic_executive_document(request: str, document_type: str, chunk_store: List[Dict[str, Any]], client: OpenAI, **kwargs: Any) -> Dict[str, Any]:
    return _generate_for_category("strategic_executive", request, document_type, chunk_store, client, **kwargs)


def generate_policy_research_document(request: str, document_type: str, chunk_store: List[Dict[str, Any]], client: OpenAI, **kwargs: Any) -> Dict[str, Any]:
    return _generate_for_category("policy_advocacy_research", request, document_type, chunk_store, client, **kwargs)


def generate_communications_document(request: str, document_type: str, chunk_store: List[Dict[str, Any]], client: OpenAI, **kwargs: Any) -> Dict[str, Any]:
    return _generate_for_category("communications_marketing", request, document_type, chunk_store, client, **kwargs)


def generate_program_document(request: str, document_type: str, chunk_store: List[Dict[str, Any]], client: OpenAI, **kwargs: Any) -> Dict[str, Any]:
    return _generate_for_category("program_service_delivery", request, document_type, chunk_store, client, **kwargs)


def generate_financial_compliance_document(request: str, document_type: str, chunk_store: List[Dict[str, Any]], client: OpenAI, **kwargs: Any) -> Dict[str, Any]:
    return _generate_for_category("financial_compliance", request, document_type, chunk_store, client, **kwargs)


def generate_event_community_document(request: str, document_type: str, chunk_store: List[Dict[str, Any]], client: OpenAI, **kwargs: Any) -> Dict[str, Any]:
    return _generate_for_category("event_community_engagement", request, document_type, chunk_store, client, **kwargs)


def _generate_for_category(category: str, request: str, document_type: str, chunk_store: List[Dict[str, Any]], client: OpenAI, **kwargs: Any) -> Dict[str, Any]:
    key = normalize_document_type(document_type)
    if DOCUMENT_CATALOG[key]["category"] != category:
        raise ValueError(f"{document_type} belongs to {DOCUMENT_CATALOG[key]['category']}, not {category}.")
    return generate_grounded_document(request, key, chunk_store, client, **kwargs)


def _plain_markdown(text: str) -> str:
    text = re.sub(r"\[([^\]]+)\]\([^\)]+\)", r"\1", text)
    text = re.sub(r"^#{1,6}\s+", "", text, flags=re.MULTILINE)
    return re.sub(r"\*\*|__|`", "", text)


def _add_docx_hyperlink(paragraph: Any, label: str, url: str) -> None:
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    relationship_id = paragraph.part.relate_to(url, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink", is_external=True)
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), relationship_id)
    run = OxmlElement("w:r")
    properties = OxmlElement("w:rPr")
    color = OxmlElement("w:color")
    color.set(qn("w:val"), "0563C1")
    underline = OxmlElement("w:u")
    underline.set(qn("w:val"), "single")
    properties.extend([color, underline])
    text = OxmlElement("w:t")
    text.text = label
    run.extend([properties, text])
    hyperlink.append(run)
    paragraph._p.append(hyperlink)


def _add_markdown_links_to_docx(paragraph: Any, text: str) -> None:
    cursor = 0
    for match in re.finditer(r"\[([^\]]+)\]\(([^\)]+)\)", text):
        paragraph.add_run(text[cursor:match.start()])
        _add_docx_hyperlink(paragraph, match.group(1), match.group(2))
        cursor = match.end()
    paragraph.add_run(text[cursor:])


def export_to_word(result: Dict[str, Any], output_path: str, title: str = "Generated Nonprofit Document") -> str:
    """Export a generated result to a styled Word document with live citations."""
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Inches, Pt, RGBColor

    document = Document()
    section = document.sections[0]
    section.top_margin = section.bottom_margin = Inches(0.8)
    section.left_margin = section.right_margin = Inches(0.85)
    normal = document.styles["Normal"]
    normal.font.name = "Arial"
    normal.font.size = Pt(10.5)
    title_paragraph = document.add_paragraph()
    title_paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
    title_run = title_paragraph.add_run(title)
    title_run.bold = True
    title_run.font.name = "Arial"
    title_run.font.size = Pt(22)
    title_run.font.color.rgb = RGBColor(31, 78, 121)

    for block in result.get("answer", "").splitlines():
        if not block.strip():
            continue
        heading = re.match(r"^(#{1,3})\s+(.+)$", block)
        if heading:
            paragraph = document.add_heading(level=min(len(heading.group(1)), 3))
            _add_markdown_links_to_docx(paragraph, heading.group(2))
        elif re.match(r"^[-*]\s+", block):
            paragraph = document.add_paragraph(style="List Bullet")
            _add_markdown_links_to_docx(paragraph, re.sub(r"^[-*]\s+", "", block))
        else:
            _add_markdown_links_to_docx(document.add_paragraph(), block)

    if result.get("citations"):
        document.add_heading("Citation Details", level=1)
        for citation in result["citations"]:
            paragraph = document.add_paragraph()
            paragraph.add_run(f"{citation['id']} - {citation['document_title']}, {citation['location']}: ").bold = True
            _add_docx_hyperlink(paragraph, "Open source document", citation.get("source_url", citation["url"]))
            excerpt = document.add_paragraph(citation["excerpt"])
            excerpt.style = document.styles["Quote"]

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    document.save(output)
    return str(output.resolve())


def _markdown_to_reportlab(text: str) -> str:
    parts: List[str] = []
    cursor = 0
    for match in re.finditer(r"\[([^\]]+)\]\(([^\)]+)\)", text):
        parts.append(html.escape(text[cursor:match.start()]))
        parts.append(f'<link href="{html.escape(match.group(2), quote=True)}" color="#0563C1"><u>{html.escape(match.group(1))}</u></link>')
        cursor = match.end()
    parts.append(html.escape(text[cursor:]))
    return "".join(parts)


def export_to_pdf(result: Dict[str, Any], output_path: str, title: str = "Generated Nonprofit Document") -> str:
    """Export a generated result to a paginated PDF with clickable citations."""
    try:
        from reportlab.lib.colors import HexColor
        from reportlab.lib.enums import TA_LEFT
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import inch
        from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer
    except ImportError as exc:
        raise RuntimeError("PDF export requires reportlab. Install it with: pip install reportlab") from exc

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="OutputTitle", parent=styles["Title"], fontName="Helvetica-Bold", fontSize=20, leading=24, textColor=HexColor("#1F4E79"), alignment=TA_LEFT, spaceAfter=18))
    styles.add(ParagraphStyle(name="OutputBody", parent=styles["BodyText"], fontName="Helvetica", fontSize=10, leading=14, spaceAfter=8))
    story: List[Any] = [Paragraph(html.escape(title), styles["OutputTitle"])]
    for block in result.get("answer", "").splitlines():
        if not block.strip():
            story.append(Spacer(1, 4))
            continue
        heading = re.match(r"^(#{1,3})\s+(.+)$", block)
        style = styles[f"Heading{min(len(heading.group(1)), 3)}"] if heading else styles["OutputBody"]
        content = heading.group(2) if heading else block
        story.append(Paragraph(_markdown_to_reportlab(content), style))
    if result.get("citations"):
        story.extend([PageBreak(), Paragraph("Citation Details", styles["Heading1"])])
        for citation in result["citations"]:
            label = f"{citation['id']} - {citation['document_title']}, {citation['location']}"
            citation_url = citation.get("source_url", citation["url"])
            story.append(Paragraph(f'<link href="{html.escape(citation_url, quote=True)}" color="#0563C1"><u>{html.escape(label)}</u></link>', styles["Heading3"]))
            story.append(Paragraph(html.escape(citation["excerpt"]), styles["OutputBody"]))
    pdf = SimpleDocTemplate(str(output), pagesize=letter, rightMargin=0.75 * inch, leftMargin=0.75 * inch, topMargin=0.7 * inch, bottomMargin=0.7 * inch, title=title)
    pdf.build(story)
    return str(output.resolve())


def export_to_powerpoint(result: Dict[str, Any], output_path: str, title: str = "Generated Nonprofit Document") -> str:
    """Export content to an editable PowerPoint deck plus citation slides."""
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    cover = presentation.slides.add_slide(presentation.slide_layouts[0])
    cover.shapes.title.text = title
    cover.placeholders[1].text = "Grounded nonprofit document output"

    clean = _plain_markdown(result.get("answer", ""))
    blocks = [block.strip() for block in re.split(r"\n(?=#{1,3}\s)|\n\s*\n", clean) if block.strip()]
    slide_chunks: List[str] = []
    for block in blocks:
        while len(block) > 1000:
            split_at = block.rfind(" ", 0, 1000)
            split_at = split_at if split_at > 400 else 1000
            slide_chunks.append(block[:split_at])
            block = block[split_at:].strip()
        if block:
            slide_chunks.append(block)
    for number, block in enumerate(slide_chunks, 1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        lines = block.splitlines()
        slide.shapes.title.text = re.sub(r"^#{1,3}\s+", "", lines[0]) if len(lines[0]) < 100 else f"Document Content {number}"
        body_text = "\n".join(lines[1:]) if len(lines) > 1 and len(lines[0]) < 100 else block
        frame = slide.placeholders[1].text_frame
        frame.clear()
        frame.word_wrap = True
        frame.paragraphs[0].text = body_text
        frame.paragraphs[0].font.size = Pt(20)
        frame.paragraphs[0].font.color.rgb = RGBColor(40, 40, 40)

    citations = result.get("citations", [])
    for start in range(0, len(citations), 5):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Citation Details"
        frame = slide.placeholders[1].text_frame
        frame.clear()
        for index, citation in enumerate(citations[start:start + 5]):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            run = paragraph.add_run()
            run.text = f"{citation['id']} - {citation['document_title']}, {citation['location']}"
            run.hyperlink.address = citation.get("source_url", citation["url"])
            run.font.color.rgb = RGBColor(5, 99, 193)
            run.font.underline = True
            paragraph.add_run().text = f"\n{citation['excerpt'][:240]}"
            paragraph.space_after = Pt(10)
            paragraph.font.size = Pt(12)

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)
    return str(output.resolve())


def export_to_excel(result: Dict[str, Any], output_path: str, title: str = "Generated Nonprofit Document") -> str:
    """Export content and a hyperlink-enabled citation register to Excel."""
    from openpyxl.styles import Alignment, Font, PatternFill

    workbook = openpyxl.Workbook()
    output_sheet = workbook.active
    output_sheet.title = "Output"
    output_sheet.append([title])
    output_sheet["A1"].font = Font(name="Arial", size=18, bold=True, color="FFFFFF")
    output_sheet["A1"].fill = PatternFill("solid", fgColor="1F4E79")
    output_sheet.column_dimensions["A"].width = 120
    for line in result.get("answer", "").splitlines():
        output_sheet.append([_plain_markdown(line)])
        output_sheet.cell(output_sheet.max_row, 1).alignment = Alignment(wrap_text=True, vertical="top")
    output_sheet.freeze_panes = "A2"

    citation_sheet = workbook.create_sheet("Citations")
    citation_sheet.append(["Citation", "Document", "Exact Location", "Excerpt", "Open Citation"])
    for cell in citation_sheet[1]:
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="1F4E79")
    for citation in result.get("citations", []):
        citation_sheet.append([citation["id"], citation["document_title"], citation["location"], citation["excerpt"], "Open exact citation"])
        link_cell = citation_sheet.cell(citation_sheet.max_row, 5)
        link_cell.hyperlink = citation.get("source_url", citation["url"])
        link_cell.style = "Hyperlink"
    widths = {"A": 12, "B": 32, "C": 28, "D": 90, "E": 22}
    for column, width in widths.items():
        citation_sheet.column_dimensions[column].width = width
    for row in citation_sheet.iter_rows(min_row=2):
        for cell in row:
            cell.alignment = Alignment(wrap_text=True, vertical="top")
    citation_sheet.freeze_panes = "A2"
    citation_sheet.auto_filter.ref = citation_sheet.dimensions

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    workbook.save(output)
    return str(output.resolve())


def export_outputs(result: Dict[str, Any], output_dir: str, formats: Sequence[str] = ("pdf", "docx", "pptx", "xlsx"), filename_stem: str = "nonprofit_output", title: str = "Generated Nonprofit Document") -> Dict[str, str]:
    """Create downloadable files and return absolute paths keyed by format."""
    exporters = {"pdf": export_to_pdf, "docx": export_to_word, "word": export_to_word, "pptx": export_to_powerpoint, "ppt": export_to_powerpoint, "xlsx": export_to_excel, "excel": export_to_excel}
    extensions = {"word": "docx", "ppt": "pptx", "excel": "xlsx"}
    exported: Dict[str, str] = {}
    for requested_format in formats:
        key = requested_format.lower().lstrip(".")
        if key not in exporters:
            raise ValueError(f"Unsupported export format: {requested_format}")
        extension = extensions.get(key, key)
        path = str(Path(output_dir) / f"{filename_stem}.{extension}")
        exported[extension] = exporters[key](result, path, title)
    return exported


def run_nonprofit_assistant(user_query: str, uploaded_file_paths: List[str], mode: str = "qa", document_type: Optional[str] = None, api_key: Optional[str] = None, chat_model: str = DEFAULT_CHAT_MODEL, embedding_model: str = DEFAULT_EMBEDDING_MODEL, citation_base_url: Optional[str] = None, export_formats: Optional[Sequence[str]] = None, export_dir: Optional[str] = None, export_filename: str = "nonprofit_output") -> Dict[str, Any]:
    """Build the repository and run grounded Q&A or document generation.

    Use mode="qa" for questions. Use mode="generate" and provide any document
    type returned by list_document_types() for drafting. Set citation_base_url
    to the app route that displays a chunk, for example
    "https://app.example.org/citations"; the stable chunk_id is appended to it.
    Set export_formats to any combination of pdf, docx, pptx, and xlsx to add
    absolute downloadable file paths under the response's "exports" key.
    """
    if not uploaded_file_paths:
        raise ValueError("At least one uploaded source document is required for grounded output.")
    client = get_openai_client(api_key)
    library = build_document_library(uploaded_file_paths)
    store = embed_chunk_store(build_chunk_store(library), client, embedding_model)
    normalized_mode = mode.lower().strip()
    if normalized_mode == "qa":
        result = answer_question(user_query, store, client, chat_model=chat_model, citation_base_url=citation_base_url)
    elif normalized_mode in {"generate", "document"}:
        if not document_type:
            raise ValueError("document_type is required when mode='generate'.")
        result = generate_grounded_document(user_query, document_type, store, client, chat_model=chat_model, citation_base_url=citation_base_url)
    else:
        # A document key may also be passed directly as mode for UI convenience.
        result = generate_grounded_document(user_query, normalized_mode, store, client, chat_model=chat_model, citation_base_url=citation_base_url)

    if export_formats:
        destination = export_dir or str(Path.cwd() / "outputs")
        title = DOCUMENT_CATALOG.get(result.get("document_type", ""), {}).get("title", "Generated Nonprofit Document")
        result["exports"] = export_outputs(result, destination, export_formats, export_filename, title)
    return result


if __name__ == "__main__":
    categories = list_document_types()
    print(f"Nonprofit RAG pipeline supports {sum(len(items) for items in categories.values())} document types.")
    for category, items in categories.items():
        print(f"- {CATEGORY_TITLES[category]}: {', '.join(items)}")
